"""Tests for the sf-04 McKinsey-grade PPTX export redesign.

Structural assertions — validates 7-slide deck, native tables, SCENE_4B
terms as searchable text, graceful degradation, adversarial Unicode
robustness, and theme coverage. Pixel-diff is out of scope for PPTX.
"""
from __future__ import annotations

import io
import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.core.exports.brand_tokens import SCENE_4B_REQUIRED_TERMS
from app.core.exports.pptx_export import generate_pptx

FIXTURE_DIR = Path(__file__).parent / "fixtures"
RUN_128 = json.loads((FIXTURE_DIR / "run_128_export.json").read_text())
RUN_NO_STATS = json.loads((FIXTURE_DIR / "run_no_stats.json").read_text())
RUN_NO_BIAS = json.loads((FIXTURE_DIR / "run_no_bias.json").read_text())
RUN_ADVERSARIAL = json.loads((FIXTURE_DIR / "run_adversarial_text.json").read_text())


def _collect_all_text(prs: Presentation) -> str:
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text_frame.text)
    return "\n".join(parts)


def _slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text_frame.text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Signature + structural smoke tests
# ---------------------------------------------------------------------------
def test_signature_returns_bytes_light():
    out = generate_pptx(RUN_128, "light")
    assert isinstance(out, bytes)
    assert len(out) > 1000


def test_signature_returns_bytes_dark():
    out = generate_pptx(RUN_128, "dark")
    assert isinstance(out, bytes)
    assert len(out) > 1000


def test_deck_has_exactly_seven_slides():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    assert len(prs.slides) == 7


def test_slide_dimensions_are_widescreen_16_9():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    # Inches(13.333) maps to 12,192,000 EMU; tolerate ±1 EMU.
    assert abs(prs.slide_width - Inches(13.333)) <= 1
    assert abs(prs.slide_height - Inches(7.5)) <= 1


# ---------------------------------------------------------------------------
# Scene 4b term coverage — deck-wide
# ---------------------------------------------------------------------------
def test_all_nine_scene4b_terms_appear_light():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    text = _collect_all_text(prs)
    missing = [t for t in SCENE_4B_REQUIRED_TERMS if t not in text]
    assert not missing, f"Missing Scene 4b terms: {missing}"


def test_all_nine_scene4b_terms_appear_dark():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "dark")))
    text = _collect_all_text(prs)
    missing = [t for t in SCENE_4B_REQUIRED_TERMS if t not in text]
    assert not missing, f"Missing Scene 4b terms: {missing}"


# ---------------------------------------------------------------------------
# Native-table assertions
# ---------------------------------------------------------------------------
def test_slide3_leaderboard_has_table():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    slide = prs.slides[2]
    assert any(sh.has_table for sh in slide.shapes), "Slide 3 must contain a native table"


def test_slide5_has_table_or_long_methodology_text():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    slide = prs.slides[4]
    has_table = any(sh.has_table for sh in slide.shapes)
    long_text = any(
        sh.has_text_frame and len(sh.text_frame.text) >= 300
        for sh in slide.shapes
    )
    assert has_table or long_text, (
        "Slide 5 must have either a pairwise matrix table or a long "
        "methodology paragraph (>=300 chars)."
    )


def test_slide4_criteria_has_table():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    slide = prs.slides[3]
    assert any(sh.has_table for sh in slide.shapes), "Slide 4 must contain a native table"


# ---------------------------------------------------------------------------
# Per-slide content assertions — per the §5 layout contract
# ---------------------------------------------------------------------------
def test_slide1_cover_contract():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    text = _slide_text(prs.slides[0])
    upper = text.upper()
    assert "MODEL EVALUATION REPORT" in upper
    assert RUN_128["run"]["name"] in text
    assert "KEY FINDING" in upper
    winner = RUN_128["models"][0]
    assert winner["name"] in text
    assert f"{winner['weighted_score']:.2f}" in text
    assert "PODIUM" in upper
    for m in RUN_128["models"][:3]:
        assert m["name"] in text
    for label in ("CANDIDATES", "PROMPTS", "JUDGES", "RUBRIC CRITERIA", "TOTAL SPEND", "MODE"):
        assert label in upper, f"scope tile {label!r} missing on cover"
    assert f"BeLLMark Run #{RUN_128['run']['id']:04d}" in text
    assert "bellmark.ai" in text


def test_slide2_executive_contract():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    text = _slide_text(prs.slides[1])
    upper = text.upper()
    winner = RUN_128["models"][0]
    assert winner["name"] in text
    assert "/ 10" in text
    for label in ("TOTAL SPEND", "TOP CLUSTER", "SAMPLE SIZE"):
        assert label in upper
    # Top-3 model names + per-prompt cost markers
    for m in RUN_128["models"][:3]:
        assert m["name"] in text
    assert text.count("/prompt") >= 3


def test_slide3_leaderboard_contract():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    slide = prs.slides[2]
    text = _slide_text(slide).lower()
    # Updated 2026-04-21 to match the 9-column contract (LC win rate → LC win,
    # plus new TOK/S and AVG LATENCY columns).
    for col in (
        "#", "model", "mean", "95% ci", "win rate",
        "tok/s", "avg latency", "lc win", "$/prompt",
    ):
        assert col in text, f"Column header {col!r} missing on leaderboard"
    # At least one row per model
    table = next(sh.table for sh in slide.shapes if sh.has_table)
    # header + models rows (allow truncation at 8 but fixture has 8)
    assert len(table.rows) == 1 + len(RUN_128["models"])


def test_slide4_criteria_contract():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    slide = prs.slides[3]
    text = _slide_text(slide)
    for c in RUN_128["run"]["criteria"]:
        assert c["name"] in text
    # Best-in-class strip present
    assert "BEST-IN-CLASS" in text.upper()


def test_slide5_stats_rigor_contract():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    text = _slide_text(prs.slides[4])
    upper = text.upper()
    for term in SCENE_4B_REQUIRED_TERMS:
        assert term in text, f"Scene 4b term {term!r} missing from slide 5"
    assert "HOW TO READ" in upper
    assert "EFFECT SIZE" in upper
    # M-code legend. Updated 2026-04-21: reasoning marker is now an icon,
    # not a (R) prefix. The legend text contains a visible-whitespace gap
    # when the model is reasoning, so check the name substring is present.
    winner_name = RUN_128["models"][0]["name"]
    assert f"M1 =" in text and winner_name in text


def test_slide6_bias_contract():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    text = _slide_text(prs.slides[5])
    upper = text.upper()
    for label in ("Position bias", "Length bias", "Self-preference", "Verbosity bias"):
        assert label in text, f"Bias panel {label!r} missing"
    assert "DETECTED" in upper or "NOT DETECTED" in upper
    assert "INTER-JUDGE AGREEMENT" in upper
    assert "FLEISS' K" in upper or "COHEN'S K" in upper
    assert "DISAGREEMENT QUESTIONS" in upper


def test_slide7_methodology_contract():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "light")))
    text = _slide_text(prs.slides[6])
    upper = text.upper()
    assert "SCOPE" in upper
    for label in ("Prompts", "Candidate models", "Judges", "Rubric criteria",
                  "Judging mode", "Judgments total", "Total spend"):
        assert label in text, f"Scope item {label!r} missing on methodology slide"
    assert "RUBRIC CRITERIA" in upper
    for c in RUN_128["run"]["criteria"]:
        assert c["name"] in text
    assert "JUDGE PANEL" in upper
    for j in RUN_128["judges"]:
        assert j["name"] in text
    assert "SIGN-OFF" in upper
    for signoff_field in ("Prepared by", "Reviewed by", "Export date", "Source", "Methodology"):
        assert signoff_field in text, f"Sign-off field {signoff_field!r} missing"


# ---------------------------------------------------------------------------
# Graceful degradation
# ---------------------------------------------------------------------------
def test_no_stats_still_produces_seven_slides_and_all_terms():
    out = generate_pptx(RUN_NO_STATS, "light")
    prs = Presentation(io.BytesIO(out))
    assert len(prs.slides) == 7
    text = _collect_all_text(prs)
    missing = [t for t in SCENE_4B_REQUIRED_TERMS if t not in text]
    assert not missing
    slide5_text = _slide_text(prs.slides[4])
    assert "unavailable" in slide5_text.lower()


def test_no_bias_still_produces_seven_slides_and_calibration():
    out = generate_pptx(RUN_NO_BIAS, "light")
    prs = Presentation(io.BytesIO(out))
    assert len(prs.slides) == 7
    slide6_text = _slide_text(prs.slides[5])
    assert "unavailable" in slide6_text.lower()
    # Calibration tiles still render
    upper = slide6_text.upper()
    assert "INTER-JUDGE AGREEMENT" in upper
    assert "DISAGREEMENT QUESTIONS" in upper


# ---------------------------------------------------------------------------
# Adversarial text — lossless per-format contract
# ---------------------------------------------------------------------------
def test_adversarial_text_does_not_raise_and_reopens():
    # This both verifies no exception AND that the output is re-openable.
    out = generate_pptx(RUN_ADVERSARIAL, "light")
    assert isinstance(out, bytes) and len(out) > 1000
    prs = Presentation(io.BytesIO(out))
    assert len(prs.slides) == 7


def test_adversarial_text_preserves_emoji_and_cjk_drops_xml_invalid():
    out = generate_pptx(RUN_ADVERSARIAL, "light")
    prs = Presentation(io.BytesIO(out))
    text = _collect_all_text(prs)
    # Emoji preserved
    assert "🔥" in text
    # CJK preserved
    assert "中文测试" in text
    # XML-invalid vertical-tab byte stripped
    assert "\x0b" not in text


# ---------------------------------------------------------------------------
# Dark theme parity
# ---------------------------------------------------------------------------
def test_dark_theme_structural_assertions_still_pass():
    prs = Presentation(io.BytesIO(generate_pptx(RUN_128, "dark")))
    assert len(prs.slides) == 7
    # Slide 3 table still there
    assert any(sh.has_table for sh in prs.slides[2].shapes)
    # All 9 terms still present
    text = _collect_all_text(prs)
    for term in SCENE_4B_REQUIRED_TERMS:
        assert term in text


# ---------------------------------------------------------------------------
# 2026-04-21 brand-parity regression tests
# ---------------------------------------------------------------------------
import hashlib
from collections import Counter
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture

from app.core.exports.brand_tokens import (
    get_tokens,
    is_reasoning_model,
)
from app.core.exports.pptx_export import _ASSETS_DIR, px_to_emu


def _png_sha(theme: str) -> str:
    path = _ASSETS_DIR / f"brain-icon-{theme}.png"
    return hashlib.sha1(path.read_bytes()).hexdigest()


def test_no_r_prefix_leaks_into_any_text_frame():
    """The `(R) ` legacy prefix must be gone from every theme + text frame."""
    for theme in ("light", "dark"):
        pptx = generate_pptx(RUN_128, theme)
        prs = Presentation(io.BytesIO(pptx))
        for sl_idx, sl in enumerate(prs.slides):
            for sh in sl.shapes:
                if sh.has_text_frame:
                    assert "(R) " not in sh.text_frame.text, (
                        f"{theme} slide {sl_idx}: '(R) ' leaked into "
                        f"{sh.text_frame.text[:120]!r}"
                    )


def test_leaderboard_uses_uppercase_nine_column_contract():
    """Slide 3 leaderboard has exactly the 9 spec-mandated UPPERCASE headers."""
    expected = [
        "#", "MODEL", "MEAN", "95% CI", "WIN RATE",
        "TOK/S", "AVG LATENCY", "LC WIN", "$/PROMPT",
    ]
    pptx = generate_pptx(RUN_128, "light")
    prs = Presentation(io.BytesIO(pptx))
    for sh in prs.slides[2].shapes:
        if isinstance(sh, GraphicFrame) and sh.has_table:
            headers = [sh.table.cell(0, i).text for i in range(len(sh.table.columns))]
            assert headers == expected, headers
            return
    raise AssertionError("no leaderboard table found on slide 3")


def test_leaderboard_tok_s_cells_render_as_number_or_dash():
    """Every data row's TOK/S cell is a '.1f' number or '—'."""
    import re as _re
    pptx = generate_pptx(RUN_128, "light")
    prs = Presentation(io.BytesIO(pptx))
    for sh in prs.slides[2].shapes:
        if isinstance(sh, GraphicFrame) and sh.has_table:
            t = sh.table
            for row_idx in range(1, len(t.rows)):
                val = t.cell(row_idx, 5).text.strip()
                assert val == "—" or _re.fullmatch(r"\d+\.\d", val), (
                    f"row {row_idx} TOK/S cell invalid: {val!r}"
                )
            return
    raise AssertionError("no leaderboard table found on slide 3")


def test_no_brain_png_pictures_emitted():
    """Brain icons were retired — the "[Reasoning …]" suffix in the
    model name is the authoritative reasoning marker across all export
    formats. No brain PNG should appear anywhere in the deck."""
    for theme in ("light", "dark"):
        png = _ASSETS_DIR / f"brain-icon-{theme}.png"
        # File may or may not still be on disk — only the SHA match matters.
        if not png.exists():
            continue
        expected_sha = _png_sha(theme)
        pptx = generate_pptx(RUN_128, theme)
        prs = Presentation(io.BytesIO(pptx))
        for sl in prs.slides:
            for sh in sl.shapes:
                if isinstance(sh, Picture):
                    assert (
                        hashlib.sha1(sh.image.blob).hexdigest() != expected_sha
                    ), f"{theme}: brain PNG leaked into deck"


def test_winner_model_cell_uses_amber_brand_color():
    """Leaderboard rank-1 cell text must match the frontend amber/yellow."""
    from app.core.exports.brand_tokens import brand
    for theme in ("light", "dark"):
        pptx = generate_pptx(RUN_128, theme)
        prs = Presentation(io.BytesIO(pptx))
        for sh in prs.slides[2].shapes:
            if isinstance(sh, GraphicFrame) and sh.has_table:
                cell = sh.table.cell(1, 1)  # row 1 = winner, col 1 = MODEL
                # Inspect first run's color via OOXML — python-pptx exposes
                # the RGB via run.font.color.rgb when set explicitly.
                run = cell.text_frame.paragraphs[0].runs[0]
                rgb = run.font.color.rgb
                expected = brand(theme, "winner_text")
                assert tuple(rgb) == tuple(expected), (
                    f"{theme}: winner cell rgb {tuple(rgb)} != "
                    f"expected {tuple(expected)}"
                )
                return
    raise AssertionError("no leaderboard table found on slide 3")


def test_runner_cells_use_green_brand_color():
    """Leaderboard rows 2..N use the frontend green brand color."""
    from app.core.exports.brand_tokens import brand
    for theme in ("light", "dark"):
        pptx = generate_pptx(RUN_128, theme)
        prs = Presentation(io.BytesIO(pptx))
        for sh in prs.slides[2].shapes:
            if isinstance(sh, GraphicFrame) and sh.has_table:
                cell = sh.table.cell(2, 1)  # row 2 = rank 2 runner
                run = cell.text_frame.paragraphs[0].runs[0]
                rgb = run.font.color.rgb
                expected = brand(theme, "runner_text")
                assert tuple(rgb) == tuple(expected), (
                    f"{theme}: runner cell rgb {tuple(rgb)} != "
                    f"expected {tuple(expected)}"
                )
                return
    raise AssertionError("no leaderboard table found on slide 3")
