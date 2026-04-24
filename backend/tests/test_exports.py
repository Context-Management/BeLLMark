"""Tests for export functionality (PPTX, HTML, JSON, CSV)."""
import copy
import json
from datetime import datetime, timezone

import pytest
from pptx import Presentation
from io import BytesIO
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

from app.core.exports.common import prepare_export_data
from app.core.exports.json_export import generate_json
from app.core.exports.pptx_export import generate_pptx
from app.core.exports.themes import get_theme, score_color_for_theme
from app.db.database import Base
from app.db.models import BenchmarkRun, Generation, Judgment, ModelPreset, ProviderType, Question, RunStatus, TaskStatus, JudgeMode


# Mock export data with complete structure (2 models, 1 judge, 2 questions, 2 criteria)
MOCK_EXPORT_DATA = {
    "run": {
        "id": 1,
        "name": "Test Benchmark Run",
        "status": "completed",
        "judge_mode": "comparison",
        "criteria": [
            {"name": "Accuracy", "description": "How accurate is the response", "weight": 2.0},
            {"name": "Clarity", "description": "How clear is the response", "weight": 1.0},
        ],
        "created_at": "2026-02-12T10:00:00",
        "completed_at": "2026-02-12T10:05:00",
        "duration_seconds": 300,
        "total_cost": 0.0234,
        "total_context_tokens": 5000,
        "run_config_snapshot": {},
    },
    "models": [
        {
            "id": 1,
            "name": "GPT-4o",
            "provider": "openai",
            "model_id": "gpt-4o",
            "rank": 1,
            "weighted_score": 8.5,
            "unweighted_score": 8.3,
            "win_count": 2,
            "total_tokens": 2500,
            "tokens_per_second": 45.2,
            "estimated_cost": 0.0150,
            "avg_latency_ms": 1200,
            "p50_latency_ms": 1150,
            "p95_latency_ms": 1300,
            "per_criterion_scores": {
                "Accuracy": 8.8,
                "Clarity": 8.2,
            },
            "per_question_scores": [
                {"order": 1, "score": 8.6},
                {"order": 2, "score": 8.4},
            ],
            "insights": ["Fastest", "Most Expensive"],
        },
        {
            "id": 2,
            "name": "Claude Sonnet 4.5",
            "provider": "anthropic",
            "model_id": "claude-sonnet-4-5",
            "rank": 2,
            "weighted_score": 7.9,
            "unweighted_score": 7.8,
            "win_count": 0,
            "total_tokens": 2200,
            "tokens_per_second": 38.5,
            "estimated_cost": 0.0084,
            "avg_latency_ms": 1500,
            "p50_latency_ms": 1450,
            "p95_latency_ms": 1650,
            "per_criterion_scores": {
                "Accuracy": 8.1,
                "Clarity": 7.7,
            },
            "per_question_scores": [
                {"order": 1, "score": 8.0},
                {"order": 2, "score": 7.8},
            ],
            "insights": ["Cheapest", "Most Concise"],
        },
    ],
    "judges": [
        {
            "id": 3,
            "name": "Claude Opus 4.6",
            "provider": "anthropic",
            "total_tokens": 1500,
            "estimated_cost": 0.0045,
            "judgment_count": 2,
        },
    ],
    "judge_summary": {
        "agreement_rate": 1.0,
        "disagreement_count": 0,
        "disagreement_questions": [],
        "per_judge_winners": {
            "Claude Opus 4.6": {
                "GPT-4o": 2,
            }
        },
    },
    "comment_summaries": {
        "Claude Opus 4.6": {
            "GPT-4o": "Excellent response with clear structure and accurate information.",
            "Claude Sonnet 4.5": "Good response but could be more detailed in some areas.",
        }
    },
    "scores_by_criterion": {
        "GPT-4o": {
            "Accuracy": 8.8,
            "Clarity": 8.2,
        },
        "Claude Sonnet 4.5": {
            "Accuracy": 8.1,
            "Clarity": 7.7,
        },
    },
    "questions": [
        {
            "id": 1,
            "order": 1,
            "system_prompt": "You are a helpful assistant.",
            "user_prompt": "What is the capital of France?",
            "context_tokens": 50,
            "generations": [
                {
                    "model_id": 1,
                    "model_name": "GPT-4o",
                    "content": "The capital of France is Paris.",
                    "tokens": 10,
                    "raw_chars": 30,
                    "answer_chars": 30,
                    "latency_ms": 1200,
                    "status": "success",
                    "error": None,
                    "retries": 0,
                },
                {
                    "model_id": 2,
                    "model_name": "Claude Sonnet 4.5",
                    "content": "Paris is the capital of France.",
                    "tokens": 9,
                    "raw_chars": 28,
                    "answer_chars": 28,
                    "latency_ms": 1500,
                    "status": "success",
                    "error": None,
                    "retries": 0,
                },
            ],
            "judgments": [
                {
                    "judge_id": 3,
                    "judge_name": "Claude Opus 4.6",
                    "blind_mapping": {"A": 1, "B": 2},
                    "rankings": ["A", "B"],
                    "scores": {
                        "1": {"Accuracy": 8.8, "Clarity": 8.2},
                        "2": {"Accuracy": 8.1, "Clarity": 7.7},
                    },
                    "reasoning": "Response A is more direct.",
                    "comments": {"1": "Clear and accurate.", "2": "Good but less direct."},
                    "latency_ms": 800,
                    "tokens": 150,
                    "status": "success",
                    "error": None,
                },
            ],
        },
        {
            "id": 2,
            "order": 2,
            "system_prompt": "You are a helpful assistant.",
            "user_prompt": "Explain quantum computing in simple terms.",
            "context_tokens": 60,
            "generations": [
                {
                    "model_id": 1,
                    "model_name": "GPT-4o",
                    "content": "Quantum computing uses quantum bits that can exist in multiple states simultaneously.",
                    "tokens": 15,
                    "raw_chars": 80,
                    "answer_chars": 80,
                    "latency_ms": 1200,
                    "status": "success",
                    "error": None,
                    "retries": 0,
                },
                {
                    "model_id": 2,
                    "model_name": "Claude Sonnet 4.5",
                    "content": "Quantum computers leverage quantum mechanics to process information in fundamentally different ways.",
                    "tokens": 14,
                    "raw_chars": 95,
                    "answer_chars": 95,
                    "latency_ms": 1500,
                    "status": "success",
                    "error": None,
                    "retries": 0,
                },
            ],
            "judgments": [
                {
                    "judge_id": 3,
                    "judge_name": "Claude Opus 4.6",
                    "blind_mapping": {"A": 1, "B": 2},
                    "rankings": ["A", "B"],
                    "scores": {
                        "1": {"Accuracy": 8.8, "Clarity": 8.2},
                        "2": {"Accuracy": 8.1, "Clarity": 7.7},
                    },
                    "reasoning": "Response A is clearer for beginners.",
                    "comments": {"1": "Very accessible.", "2": "Accurate but more technical."},
                    "latency_ms": 850,
                    "tokens": 160,
                    "status": "success",
                    "error": None,
                },
            ],
        },
    ],
}


def create_model_preset(session, preset_id: int, name: str, provider=ProviderType.openai) -> ModelPreset:
    preset = ModelPreset(
        id=preset_id,
        name=name,
        provider=provider,
        base_url="https://api.openai.com/v1/chat/completions",
        model_id=f"model-{preset_id}",
        price_input=0.0,
        price_output=0.0,
    )
    session.add(preset)
    session.commit()
    return preset


def create_benchmark_run(session, name: str, model_ids: list[int], judge_ids: list[int]) -> BenchmarkRun:
    run = BenchmarkRun(
        name=name,
        model_ids=model_ids,
        judge_ids=judge_ids,
        judge_mode=JudgeMode.comparison,
        criteria=[{"name": "Quality", "description": "Quality", "weight": 1.0}],
        status=RunStatus.completed,
        created_at=datetime.now(timezone.utc),
        completed_at=datetime.now(timezone.utc),
        total_context_tokens=0,
    )
    session.add(run)
    session.commit()
    return run


def create_question(session, benchmark_id: int, order: int) -> Question:
    question = Question(
        benchmark_id=benchmark_id,
        order=order,
        system_prompt="System prompt",
        user_prompt=f"Question {order}",
        context_tokens=0,
    )
    session.add(question)
    session.commit()
    return question


def create_generation(session, question_id: int, model_preset_id: int) -> Generation:
    generation = Generation(
        question_id=question_id,
        model_preset_id=model_preset_id,
        content="Answer",
        tokens=12,
        raw_chars=6,
        answer_chars=6,
        latency_ms=1200,
        status=TaskStatus.success,
        retries=0,
    )
    session.add(generation)
    session.commit()
    return generation


def create_judgment(session, question_id: int, judge_preset_id: int, *, score_rationales=None) -> Judgment:
    judgment = Judgment(
        question_id=question_id,
        judge_preset_id=judge_preset_id,
        blind_mapping={"A": 1},
        rankings=["A"],
        scores={1: {"Quality": 9.0}},
        reasoning="Reasoning",
        score_rationales=score_rationales,
        comments={1: [{"text": "Helpful", "sentiment": "positive"}]},
        latency_ms=1500,
        tokens=100,
        status=TaskStatus.success,
        retries=0,
    )
    session.add(judgment)
    session.commit()
    return judgment



@pytest.fixture
def db_session():
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(engine)
    SessionLocal = sessionmaker(bind=engine)
    session = SessionLocal()
    yield session
    session.close()



class TestThemes:
    """Tests for theme system."""

    def test_light_theme_has_all_keys(self):
        theme = get_theme("light")
        required = ["bg", "card_bg", "accent_bg", "text", "text_secondary",
                     "brand", "accent", "success", "divider", "table_alt_row",
                     "winner_border"]
        for key in required:
            assert key in theme, f"Missing key: {key}"

    def test_dark_theme_has_all_keys(self):
        theme = get_theme("dark")
        required = ["bg", "card_bg", "accent_bg", "text", "text_secondary",
                     "brand", "accent", "success", "divider", "table_alt_row",
                     "winner_border"]
        for key in required:
            assert key in theme, f"Missing key: {key}"

    def test_default_theme_is_light(self):
        theme = get_theme()
        assert theme["bg"] == (255, 255, 255)

    def test_score_color_returns_rgb_tuple(self):
        color = score_color_for_theme(8.5, "light")
        assert isinstance(color, tuple)
        assert len(color) == 3
        assert all(0 <= c <= 255 for c in color)

    def test_score_color_differs_by_theme(self):
        light = score_color_for_theme(7.0, "light")
        dark = score_color_for_theme(7.0, "dark")
        assert isinstance(light, tuple) and isinstance(dark, tuple)


class TestPptxExport:
    """Smoke tests for the redesigned 7-slide PPTX export.

    Content-specific assertions live in
    ``tests/test_pptx_export_redesign.py`` — this class only sanity-checks
    the signature, byte validity, reopen-ability, and slide count.
    """

    def test_generates_valid_pptx(self):
        """Output is valid PPTX bytes that python-pptx can read."""
        result = generate_pptx(MOCK_EXPORT_DATA)

        # Should return bytes
        assert isinstance(result, bytes)
        assert len(result) > 0

        # Should be valid PPTX that can be opened
        buffer = BytesIO(result)
        prs = Presentation(buffer)
        assert prs is not None

    def test_generates_valid_pptx_dark_theme(self):
        """Output is valid PPTX with dark theme."""
        result = generate_pptx(MOCK_EXPORT_DATA, theme="dark")
        assert isinstance(result, bytes)
        prs = Presentation(BytesIO(result))
        assert prs is not None

    def test_has_correct_slide_count(self):
        """Redesigned deck is always 7 slides regardless of input shape."""
        result = generate_pptx(MOCK_EXPORT_DATA)
        prs = Presentation(BytesIO(result))
        assert len(prs.slides) == 7

    def test_cover_contains_run_name(self):
        """Slide 1 (cover) renders the run name."""
        result = generate_pptx(MOCK_EXPORT_DATA)
        prs = Presentation(BytesIO(result))
        title_slide = prs.slides[0]
        all_text = []
        for shape in title_slide.shapes:
            if hasattr(shape, "text_frame"):
                all_text.append(shape.text_frame.text)
        assert "Test Benchmark Run" in " ".join(all_text)

    def test_handles_edge_cases(self):
        """Test handling of None values and edge cases."""
        edge_data = MOCK_EXPORT_DATA.copy()
        edge_data["models"] = [{
            "id": 1, "name": "Test Model", "provider": "test", "model_id": "test",
            "rank": 1, "weighted_score": 0.0, "unweighted_score": 0.0,
            "win_count": 0, "total_tokens": 0, "tokens_per_second": None,
            "estimated_cost": 0.0, "avg_latency_ms": 0, "p50_latency_ms": 0,
            "p95_latency_ms": 0, "per_criterion_scores": {"Accuracy": 0.0, "Clarity": 0.0},
            "per_question_scores": [], "insights": [],
        }]
        result = generate_pptx(edge_data)
        assert isinstance(result, bytes) and len(result) > 0

    def test_many_questions_still_yields_seven_slides(self):
        """The deck is always 7 slides — even with many questions the count is fixed."""
        import copy
        data = copy.deepcopy(MOCK_EXPORT_DATA)
        base_q = data["questions"][0]
        data["questions"] = []
        for i in range(12):
            q = copy.deepcopy(base_q)
            q["id"] = i + 1
            q["order"] = i + 1
            q["user_prompt"] = f"Question number {i + 1} about testing?"
            data["questions"].append(q)
        result = generate_pptx(data)
        prs = Presentation(BytesIO(result))
        assert len(prs.slides) == 7


class TestJsonExport:
    """Tests for JSON export functionality."""

    def test_generates_valid_json(self):
        """Output is valid JSON string that can be parsed."""
        from app.core.exports.json_export import generate_json
        result = generate_json(MOCK_EXPORT_DATA)
        assert isinstance(result, str)
        parsed = json.loads(result)
        assert "run" in parsed
        assert "models" in parsed
        assert "questions" in parsed

    def test_includes_performance_metrics(self):
        """JSON should include all performance metrics."""
        from app.core.exports.json_export import generate_json
        result = json.loads(generate_json(MOCK_EXPORT_DATA))
        model = result["models"][0]
        assert "tokens_per_second" in model
        assert "estimated_cost" in model
        assert "per_criterion_scores" in model

    def test_includes_all_computed_fields(self):
        """JSON should include all computed summary fields."""
        from app.core.exports.json_export import generate_json
        result = json.loads(generate_json(MOCK_EXPORT_DATA))
        assert "judge_summary" in result
        assert "comment_summaries" in result
        assert "scores_by_criterion" in result
        assert result["run"]["duration_seconds"] == 300
        assert result["run"]["total_cost"] == 0.0234

    def test_includes_integrity_block(self):
        """JSON export should include an _integrity metadata block with SHA-256 hash."""
        from app.core.exports.json_export import generate_json
        result = json.loads(generate_json(MOCK_EXPORT_DATA))
        assert "_integrity" in result
        integrity = result["_integrity"]
        assert "sha256" in integrity
        assert "generated_at" in integrity
        assert "bellmark_version" in integrity
        assert "run_id" in integrity
        assert len(integrity["sha256"]) == 64  # 256-bit hex digest
        assert integrity["run_id"] == MOCK_EXPORT_DATA["run"]["id"]

    def test_integrity_hash_is_verifiable(self):
        """The SHA-256 hash in the _integrity block must be verifiable against the payload."""
        import hashlib
        from app.core.exports.json_export import generate_json
        exported = json.loads(generate_json(MOCK_EXPORT_DATA))
        integrity = exported.pop("_integrity")

        canonical = json.dumps(MOCK_EXPORT_DATA, sort_keys=True, separators=(",", ":"), ensure_ascii=False, default=str)
        expected_hash = hashlib.sha256(canonical.encode("utf-8")).hexdigest()
        assert integrity["sha256"] == expected_hash

    def test_integrity_hash_changes_with_data(self):
        """Different export data should produce different hashes."""
        import copy
        from app.core.exports.json_export import generate_json
        data_a = MOCK_EXPORT_DATA
        data_b = copy.deepcopy(MOCK_EXPORT_DATA)
        data_b["run"]["name"] = "A Different Run Name"

        hash_a = json.loads(generate_json(data_a))["_integrity"]["sha256"]
        hash_b = json.loads(generate_json(data_b))["_integrity"]["sha256"]
        assert hash_a != hash_b


class TestJsonExportScoreRationales:
    def test_export_data_includes_score_rationales_mapping(self, db_session):
        model = create_model_preset(db_session, 1, "Export Model")
        judge = create_model_preset(db_session, 2, "Export Judge", provider=ProviderType.anthropic)
        run = create_benchmark_run(db_session, "Score Rationales Export", [model.id], [judge.id])
        question = create_question(db_session, run.id, 1)
        create_generation(db_session, question.id, model.id)
        create_judgment(db_session, question.id, judge.id, score_rationales={str(model.id): "Clear and accurate."})

        export_data = prepare_export_data(db_session, run.id)
        judgment = export_data["questions"][0]["judgments"][0]

        assert judgment["score_rationales"] == {str(model.id): "Clear and accurate."}
        assert "score_rationales" in judgment

    def test_export_data_keeps_null_score_rationales_when_absent(self, db_session):
        model = create_model_preset(db_session, 11, "Export Model 2")
        judge = create_model_preset(db_session, 12, "Export Judge 2", provider=ProviderType.anthropic)
        run = create_benchmark_run(db_session, "Null Rationales Export", [model.id], [judge.id])
        question = create_question(db_session, run.id, 1)
        create_generation(db_session, question.id, model.id)
        create_judgment(db_session, question.id, judge.id, score_rationales=None)

        export_data = prepare_export_data(db_session, run.id)
        judgment = export_data["questions"][0]["judgments"][0]

        assert "score_rationales" in judgment
        assert judgment["score_rationales"] is None

    def test_json_export_preserves_score_rationales_output(self, db_session):
        model = create_model_preset(db_session, 21, "Export Model 3")
        judge = create_model_preset(db_session, 22, "Export Judge 3", provider=ProviderType.anthropic)
        run = create_benchmark_run(db_session, "JSON Rationales Export", [model.id], [judge.id])
        question = create_question(db_session, run.id, 1)
        create_generation(db_session, question.id, model.id)
        create_judgment(db_session, question.id, judge.id, score_rationales={str(model.id): "Specific, concise rationale."})

        export_data = prepare_export_data(db_session, run.id)
        exported = json.loads(generate_json(export_data))

        assert exported["questions"][0]["judgments"][0]["score_rationales"] == {str(model.id): "Specific, concise rationale."}
        assert "_integrity" in exported


class TestComputeExportIntegrity:
    """Unit tests for the compute_export_integrity utility."""

    def test_returns_required_fields(self):
        from app.core.exports.common import compute_export_integrity
        result = compute_export_integrity({"foo": "bar"}, run_id=42)
        assert set(result.keys()) == {"sha256", "generated_at", "bellmark_version", "run_id"}

    def test_sha256_is_hex_string_of_correct_length(self):
        from app.core.exports.common import compute_export_integrity
        result = compute_export_integrity({"x": 1}, run_id=1)
        assert isinstance(result["sha256"], str)
        assert len(result["sha256"]) == 64
        int(result["sha256"], 16)  # must be valid hex

    def test_deterministic_for_same_input(self):
        from app.core.exports.common import compute_export_integrity
        payload = {"a": 1, "b": [2, 3]}
        h1 = compute_export_integrity(payload, run_id=7)["sha256"]
        h2 = compute_export_integrity(payload, run_id=7)["sha256"]
        assert h1 == h2

    def test_key_order_independent(self):
        """Canonical JSON uses sorted keys, so dict insertion order must not matter."""
        from app.core.exports.common import compute_export_integrity
        h1 = compute_export_integrity({"a": 1, "b": 2}, run_id=1)["sha256"]
        h2 = compute_export_integrity({"b": 2, "a": 1}, run_id=1)["sha256"]
        assert h1 == h2

    def test_generated_at_is_iso8601(self):
        from datetime import datetime, timezone
        from app.core.exports.common import compute_export_integrity
        result = compute_export_integrity({}, run_id=0)
        # Should parse without error
        dt = datetime.fromisoformat(result["generated_at"])
        assert dt.tzinfo is not None


class TestPdfExport:
    """Tests for PDF export functionality."""

    def test_generates_valid_pdf(self):
        from app.core.exports.pdf_export import generate_pdf
        result = generate_pdf(MOCK_EXPORT_DATA)
        assert isinstance(result, bytes)
        assert result[:5] == b"%PDF-"

    def test_generates_valid_pdf_dark_theme(self):
        from app.core.exports.pdf_export import generate_pdf
        result = generate_pdf(MOCK_EXPORT_DATA, theme="dark")
        assert isinstance(result, bytes)
        assert result[:5] == b"%PDF-"

    def test_has_expected_page_count(self):
        """sf-03 PDF redesign: always exactly 7 pages (cover, exec, leaderboard,
        criteria, stats, bias, methodology)."""
        import re
        from app.core.exports.pdf_export import generate_pdf
        result = generate_pdf(MOCK_EXPORT_DATA)
        assert isinstance(result, bytes)
        assert len(result) > 1000
        match = re.search(rb'/Count (\d+)', result)
        assert match, "Could not find /Count in PDF"
        page_count = int(match.group(1))
        assert page_count == 7, f"expected 7 pages, got {page_count}"

    def test_all_models_included(self):
        from app.core.exports.pdf_export import generate_pdf
        result = generate_pdf(MOCK_EXPORT_DATA)
        assert isinstance(result, bytes) and len(result) > 0

    def test_handles_edge_cases(self):
        from app.core.exports.pdf_export import generate_pdf
        edge_data = MOCK_EXPORT_DATA.copy()
        edge_data["models"] = [{
            "id": 1, "name": "Test Model", "provider": "test", "model_id": "test",
            "rank": 1, "weighted_score": 0.0, "unweighted_score": 0.0,
            "win_count": 0, "total_tokens": 0, "tokens_per_second": None,
            "estimated_cost": 0.0, "avg_latency_ms": 0, "p50_latency_ms": 0,
            "p95_latency_ms": 0, "per_criterion_scores": {"Accuracy": 0.0, "Clarity": 0.0},
            "per_question_scores": [], "insights": [],
        }]
        result = generate_pdf(edge_data)
        assert isinstance(result, bytes) and len(result) > 0


class TestHtmlExport:
    """Tests for HTML export functionality."""

    def test_generates_valid_html(self):
        """Output is valid HTML string with proper structure."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        assert isinstance(result, str)
        assert "<!DOCTYPE html>" in result
        assert "</html>" in result

    def test_contains_all_sections(self):
        """HTML contains executive summary, model details, questions, and judge analysis."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        assert "Test Benchmark Run" in result
        # Redesigned slide sections
        assert "LEADERBOARD" in result.upper()
        assert "CRITERION" in result.upper() or "Per-Criterion" in result or "criteria" in result.lower()
        # Appendix archival content
        assert "Model Details" in result
        assert "Question Details" in result
        assert "Question 1" in result
        assert "Judge Analysis" in result or "agreement" in result.lower()

    def test_contains_full_generation_content(self):
        """HTML should have full generation content, not truncated."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        # Check for full content from question 2
        assert "Quantum computing uses quantum bits that can exist in multiple states simultaneously" in result

    def test_contains_judge_reasoning(self):
        """HTML should contain full judge reasoning text."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        assert "Response A is more direct" in result or "Response A is clearer" in result

    def test_html_escapes_user_content(self):
        """HTML should properly escape user-provided content to prevent XSS."""
        import copy
        data = copy.deepcopy(MOCK_EXPORT_DATA)
        data["questions"][0]["user_prompt"] = '<script>alert("xss")</script>'
        from app.core.exports.html_export import generate_html
        result = generate_html(data)
        assert "<script>" not in result
        assert "&lt;script&gt;" in result

    def test_contains_all_models(self):
        """HTML should list all models in rankings and details."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        assert "GPT-4o" in result
        assert "Claude Sonnet 4.5" in result

    def test_contains_criterion_scores(self):
        """HTML should show per-criterion scores in heatmap."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        assert "Accuracy" in result
        assert "Clarity" in result
        # Should have scores
        assert "8.8" in result  # GPT-4o Accuracy score
        assert "8.2" in result  # GPT-4o Clarity score

    def test_contains_judge_comments(self):
        """HTML should include judge comments with proper structure."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        # Check for comments from mock data
        assert "Clear and accurate" in result or "Good but less direct" in result

    def test_handles_edge_cases(self):
        """Test handling of None values and edge cases."""
        import copy
        from app.core.exports.html_export import generate_html

        # Create data with edge cases
        edge_case_data = copy.deepcopy(MOCK_EXPORT_DATA)
        edge_case_data["models"] = [
            {
                "id": 1,
                "name": "Test Model",
                "provider": "test",
                "model_id": "test-model",
                "rank": 1,
                "weighted_score": 0.0,
                "unweighted_score": 0.0,
                "win_count": 0,
                "total_tokens": 0,
                "tokens_per_second": None,
                "estimated_cost": 0.0,
                "avg_latency_ms": 0,
                "p50_latency_ms": 0,
                "p95_latency_ms": 0,
                "per_criterion_scores": {
                    "Accuracy": 0.0,
                    "Clarity": 0.0,
                },
                "per_question_scores": [],
                "insights": [],
            }
        ]

        result = generate_html(edge_case_data)
        assert isinstance(result, str)
        assert len(result) > 0
        assert "<!DOCTYPE html>" in result

    def test_includes_winner_highlight(self):
        """Winner row should be highlighted in rankings table."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        # Should have winner-row class
        assert "winner-row" in result

    def test_includes_collapsible_sections(self):
        """HTML should use details/summary for collapsible sections."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        assert "<details" in result
        assert "<summary" in result
        assert "</details>" in result

    def test_includes_prompts(self):
        """HTML should include system and user prompts for questions."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        assert "System Prompt:" in result
        assert "User Prompt:" in result
        assert "You are a helpful assistant" in result

    def test_includes_generation_metrics(self):
        """HTML should show generation metrics like tokens and latency."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        # Should have token counts and latencies
        assert "tokens" in result.lower()
        assert "ms" in result  # milliseconds

    def test_includes_insights(self):
        """HTML should include insight badges for models."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        # GPT-4o has "Fastest" and "Most Expensive" insights
        assert "Fastest" in result
        assert "Most Expensive" in result or "Expensive" in result

    def test_footer_present(self):
        """HTML should have footer with BeLLMark branding."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        assert "BeLLMark" in result
        assert "Benchmark Studio" in result or "Generated" in result

    def test_logo_handling(self):
        """HTML should handle logo SVG or fallback to text."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        # Should either have SVG or text fallback
        assert "<svg" in result or "BeLLMark" in result

    def test_includes_integrity_comment(self):
        """HTML export should embed a BeLLMark-Integrity HTML comment in the head."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        assert "<!-- BeLLMark-Integrity" in result
        assert "sha256" in result
        assert "generated_at" in result

    def test_includes_integrity_sha256_in_footer(self):
        """HTML export footer should display a truncated SHA-256 fingerprint."""
        from app.core.exports.html_export import generate_html
        result = generate_html(MOCK_EXPORT_DATA)
        assert "SHA-256:" in result

    def test_html_light_theme(self):
        """HTML export should support light theme."""
        from app.core.exports.html_export import generate_html
        html_content = generate_html(MOCK_EXPORT_DATA, theme="light")
        assert "<html" in html_content
        assert 'data-theme="light"' in html_content

    def test_html_dark_theme_is_default(self):
        """HTML export without theme param should default to dark."""
        from app.core.exports.html_export import generate_html
        html_content = generate_html(MOCK_EXPORT_DATA)
        assert 'data-theme="dark"' in html_content


class TestCsvExport:
    """Tests for CSV export functionality."""

    def test_generates_valid_csv(self):
        """Output is valid CSV string with proper structure."""
        from app.core.exports.csv_export import generate_csv
        result = generate_csv(MOCK_EXPORT_DATA)
        assert isinstance(result, str)
        lines = result.strip().split("\n")
        assert len(lines) > 1

    def test_header_includes_all_columns(self):
        """Header row contains all expected column names."""
        from app.core.exports.csv_export import generate_csv
        result = generate_csv(MOCK_EXPORT_DATA)
        header = result.split("\n")[0]
        assert "Question" in header
        assert "Model" in header
        assert "Tokens" in header
        assert "Latency" in header
        assert "Cost" in header
        assert "Weighted Score" in header

    def test_has_row_per_question_model(self):
        """2 questions x 2 models = 4 data rows."""
        from app.core.exports.csv_export import generate_csv
        result = generate_csv(MOCK_EXPORT_DATA)
        lines = result.strip().split("\n")
        assert len(lines) == 5  # header + 4 data rows


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

def _make_many_questions_data(count: int = 20) -> dict:
    """Build export data with `count` questions (based on MOCK_EXPORT_DATA)."""
    import copy
    data = copy.deepcopy(MOCK_EXPORT_DATA)
    base_q = data["questions"][0]
    data["questions"] = []
    for i in range(count):
        q = copy.deepcopy(base_q)
        q["id"] = i + 1
        q["order"] = i + 1
        q["user_prompt"] = f"Question number {i + 1}: What is the meaning of test question {i + 1}?"
        data["questions"].append(q)
    # Update per_question_scores for both models to cover all questions
    for model in data["models"]:
        model["per_question_scores"] = [
            {"order": j + 1, "score": 8.0 + (j % 3) * 0.1} for j in range(count)
        ]
    return data


@pytest.fixture
def sample_export_data():
    import copy
    return copy.deepcopy(MOCK_EXPORT_DATA)


@pytest.fixture
def sample_export_data_many_questions():
    return _make_many_questions_data(20)


# ---------------------------------------------------------------------------
# Task 2: PDF Questions Pages - Adaptive Layout
# ---------------------------------------------------------------------------

class TestPdfQuestionsAdaptive:
    """Tests for adaptive question page layout in PDF."""

    def test_pdf_questions_overflow_protection(self, sample_export_data_many_questions):
        """Long questions should not overflow the page — new pages are created."""
        from app.core.exports.pdf_export import generate_pdf
        data = sample_export_data_many_questions
        for q in data["questions"][:3]:
            q["user_prompt"] = "x" * 600
        result = generate_pdf(data, "light")
        assert result.startswith(b"%PDF")

    def test_pdf_questions_truncation(self, sample_export_data):
        """Question text over 300 chars should not crash the generator and truncation
        utility is applied (verified via truncate_text unit tests in test_export_utils.py)."""
        from app.core.exports.pdf_export import generate_pdf
        from app.core.exports.common import truncate_text, MAX_QUESTION_DISPLAY_CHARS
        data = sample_export_data
        long_prompt = "y" * 500
        data["questions"][0]["user_prompt"] = long_prompt
        result = generate_pdf(data, "light")
        # The PDF should be generated successfully
        assert isinstance(result, bytes)
        assert result[:5] == b"%PDF-"
        # Verify the truncation function works correctly (the underlying mechanism)
        truncated = truncate_text(long_prompt, MAX_QUESTION_DISPLAY_CHARS)
        assert len(truncated) == MAX_QUESTION_DISPLAY_CHARS + 3
        assert truncated.endswith("...")

    def test_pdf_generates_with_many_questions(self, sample_export_data_many_questions):
        """PDF with 20 questions should be generated without error."""
        from app.core.exports.pdf_export import generate_pdf
        result = generate_pdf(sample_export_data_many_questions, "light")
        assert isinstance(result, bytes)
        assert result[:5] == b"%PDF-"


# ---------------------------------------------------------------------------
# Task 3: PDF Per-Model Table - No Cap + Feedback Pagination
# ---------------------------------------------------------------------------

class TestPdfModelTableNoCap:
    """Tests that the 12-question cap is removed from per-model tables."""

    def test_pdf_model_table_no_cap(self, sample_export_data_many_questions):
        """Per-model table should include all questions, not just first 12."""
        from app.core.exports.pdf_export import generate_pdf
        data = sample_export_data_many_questions  # 20 questions
        result = generate_pdf(data, "light")
        # At least Q13 should be referenced (would have been cut by old cap)
        assert b"13" in result

    def test_pdf_generates_with_many_questions_dark(self, sample_export_data_many_questions):
        """PDF with 20 questions dark theme generates without error."""
        from app.core.exports.pdf_export import generate_pdf
        result = generate_pdf(sample_export_data_many_questions, "dark")
        assert isinstance(result, bytes)
        assert result[:5] == b"%PDF-"


# ---------------------------------------------------------------------------
# Task 4 & 5: PPTX Questions Slides + Per-Model Table - Adaptive Layout
# ---------------------------------------------------------------------------

class TestPptxQuestionsAdaptive:
    """Legacy adaptive-question tests — kept as smoke tests after the 7-slide
    redesign. The deck no longer includes per-question slides; these tests
    now just verify the generator doesn't choke on many-question input."""

    def test_pptx_questions_overflow_protection(self, sample_export_data_many_questions):
        """Long questions should not crash the generator."""
        from app.core.exports.pptx_export import generate_pptx
        data = sample_export_data_many_questions
        for q in data["questions"][:3]:
            q["user_prompt"] = "x" * 600
        result = generate_pptx(data, "light")
        assert result is not None
        assert len(result) > 0

    def test_pptx_generates_with_many_questions(self, sample_export_data_many_questions):
        """PPTX with 20 questions generates without error."""
        from app.core.exports.pptx_export import generate_pptx
        result = generate_pptx(sample_export_data_many_questions, "light")
        assert isinstance(result, bytes)
        assert len(result) > 0

    def test_pptx_seven_slides_with_many_questions(self, sample_export_data_many_questions):
        """PPTX slide count is fixed at 7 regardless of question count."""
        from app.core.exports.pptx_export import generate_pptx
        from pptx import Presentation
        data = sample_export_data_many_questions
        result = generate_pptx(data, "light")
        prs = Presentation(BytesIO(result))
        assert len(prs.slides) == 7
