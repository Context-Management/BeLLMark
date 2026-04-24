"""PPTX export generator for BeLLMark benchmark results — McKinsey-grade redesign.

Rewrites the legacy 11-slide export as a 7-slide deck (cover, executive,
leaderboard, criteria, stats-rigor, bias, methodology) using native
python-pptx shape primitives: text boxes, rectangles, and tables. All
colours come from :mod:`app.core.exports.brand_tokens`, so light/dark
themes stay in lockstep with the frontend.

Canonical slide geometry: 13.333″ × 7.5″ (PowerPoint 16:9). Layout
coordinates are expressed in a 1920×1080 logical pixel grid and converted
to EMU via a single scale factor (``6350``) so both axes share the same
mapping.

Public API is UNCHANGED::

    def generate_pptx(data: dict, theme: str = "light") -> bytes: ...

The output is always a 7-slide deck regardless of whether the optional
``statistics`` / ``bias_report`` / ``calibration_report`` analytics blocks
are present — missing blocks yield placeholder text on slides 5/6.
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from app.core.exports.brand_tokens import (
    SCENE_4B_REQUIRED_TERMS,
    BrandTokens,
    _fmt_latency_ms,
    brand,
    get_tokens,
    is_reasoning_model,
    sanitize_text,
)
from app.core.exports.common import format_cost

# ---------------------------------------------------------------------------
# Geometry helpers — single source of truth for logical-px → EMU conversion.
# ---------------------------------------------------------------------------
_SCALE = 6350  # 12_192_000 EMU / 1920 logical px (= 6_858_000 / 1080)


def lx(px: float) -> Emu:
    """Convert a logical x-coordinate (0..1920 px) to EMU."""
    return Emu(int(round(px * _SCALE)))


def ly(py: float) -> Emu:
    """Convert a logical y-coordinate (0..1080 px) to EMU."""
    return Emu(int(round(py * _SCALE)))


# 96-DPI pixel → EMU converter. Used ONLY for picture sizing and the small
# offsets used when placing a brain icon relative to an already-anchored
# text box/table cell. The deck's 1920×1080 grid continues to use lx/ly.
_PPTX_PX_TO_EMU = 9525


def px_to_emu(px: float) -> Emu:
    """96-DPI pixel → EMU (for picture sizing and brain-icon offsets only)."""
    return Emu(int(round(px * _PPTX_PX_TO_EMU)))


_ASSETS_DIR = Path(__file__).parent / "assets"


def _add_brain_picture(
    slide,
    *,
    x_emu,
    y_emu,
    size_px: int,
    theme: str,
) -> None:
    """Deprecated no-op. The reasoning marker is now the "[Reasoning …]"
    suffix already in the model name — a brain icon alongside it was
    redundant. Kept as a stub so legacy tests calling it directly don't
    explode; new code should not call this."""
    return None


def _maybe_add_brain(
    slide,
    *,
    entity: dict,
    snapshots: list | None,
    theme: str,
    anchor_x_px: float,
    anchor_y_px: float,
    icon_x_off_px: int,
    icon_y_off_px: int,
    size_px: int,
) -> None:
    """Deprecated no-op. See :func:`_add_brain_picture`."""
    return None
    # unreachable legacy code kept for reference only
    _add_brain_picture(  # pragma: no cover
        slide,
        x_emu=lx(anchor_x_px) + px_to_emu(icon_x_off_px),
        y_emu=ly(anchor_y_px) + px_to_emu(icon_y_off_px),
        size_px=size_px,
        theme=theme,
    )


# ---------------------------------------------------------------------------
# Static methodology paragraph — contains every SCENE_4B_REQUIRED_TERM verbatim.
# ---------------------------------------------------------------------------
# Deliberately wordy: Scene 4b's camera pan across the stats-rigor slide needs
# every term as real, searchable text. Kept as a module constant so test fixtures
# and human readers see the same ASCII paragraph regardless of run contents.
_METHODOLOGY_PARAGRAPH = (
    "BeLLMark reports win rates with a Wilson 95% CI computed on the "
    "per-question winner counts. Score differences between models are "
    "summarised with Bootstrap confidence intervals over judge-level "
    "means. Paired comparisons run a Wilcoxon signed-rank test on the "
    "per-question score vectors, with p-values adjusted using the "
    "Holm-Bonferroni correction for multiple comparisons. Effect sizes "
    "use Cohen's d, thresholded as negligible / small / medium / large. "
    "Bias screens run before ranking: Position bias from presentation-"
    "order correlation, Length bias from Spearman correlation between "
    "response length and score, Self-preference from same-provider mean "
    "score differences, and Verbosity bias from judge-reasoning length "
    "correlation with awarded scores."
)

# Sanity check at import time — cheap and never fires in practice but guards
# against accidental edits that break the Scene 4b contract.
for _term in SCENE_4B_REQUIRED_TERMS:
    assert _term in _METHODOLOGY_PARAGRAPH, (
        f"Static methodology paragraph must contain {_term!r} verbatim."
    )


# ---------------------------------------------------------------------------
# Low-level shape primitives
# ---------------------------------------------------------------------------
def _rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(color[0], color[1], color[2])


def _set_background(slide, tokens: BrandTokens) -> None:
    """Add a full-bleed background rectangle filled with ``tokens['background']``."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, lx(0), ly(0), lx(1920), ly(1080)
    )
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb(tokens["background"])
    # Move to back so all subsequent shapes render on top.
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)


def _add_text(
    slide,
    text: str,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    color: tuple[int, int, int],
    size: int = 14,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    """Drop an editable text box at (x, y) with explicit typography."""
    box = slide.shapes.add_textbox(lx(x), ly(y), lx(w), ly(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = sanitize_text(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _add_rect(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: tuple[int, int, int],
    border: tuple[int, int, int] | None = None,
    border_width_pt: float = 0.75,
):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx(x), ly(y), lx(w), ly(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb(fill)
    if border is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = _rgb(border)
        rect.line.width = Pt(border_width_pt)
    return rect


def _add_stat_tile(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    value: str,
    tokens: BrandTokens,
):
    """Scope/stat tile — label on top, value prominent below."""
    _add_rect(
        slide, x=x, y=y, w=w, h=h,
        fill=tokens["card"], border=tokens["border"],
    )
    _add_text(
        slide, label.upper(),
        x=x + 20, y=y + 18, w=w - 40, h=28,
        color=tokens["muted_foreground"], size=11, bold=True,
    )
    # Auto-shrink the value font for long strings like "COMPARISON" so they
    # don't wrap inside a narrow (w=280) scope tile. 22pt @ 280/6 tiles = the
    # max that fits ~7 uppercase chars; longer values step down to 16pt.
    value_size = 22 if len(value) <= 7 else 16
    _add_text(
        slide, value,
        x=x + 20, y=y + 52, w=w - 40, h=h - 60,
        color=tokens["foreground"], size=value_size, bold=True,
    )


def _reset_table_style(table) -> None:
    """Neutralise PowerPoint's default banded table style.

    By default, python-pptx ``add_table`` creates a table with the
    "Medium Style 2 - Accent 1" template applied, which ships with
    alternating white/cream row fills and its own text colors. On a dark
    theme slide this produces unreadable light-text-on-white-stripe rows.

    We zero every style flag and strip the ``<a:tableStyleId>`` reference
    from the underlying XML so the cells inherit NOTHING from the
    template — everything (fill, text color) must be set explicitly on
    each cell by the caller.
    """
    table.first_row = False
    table.first_col = False
    table.last_row = False
    table.last_col = False
    table.horz_banding = False
    table.vert_banding = False
    # Remove the table style id element — this is what actually carries the
    # "Medium Style 2 - Accent 1" default template in the OOXML. Without
    # this removal, horz_banding=False is not enough on some viewers.
    tbl = table._tbl
    for child in list(tbl.iter()):
        if child.tag.endswith("}tableStyleId"):
            child.getparent().remove(child)


def _set_cell(
    cell,
    text: str,
    *,
    color: tuple[int, int, int],
    size: int = 12,
    bold: bool = False,
    fill: tuple[int, int, int] | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    cell.text = ""  # reset
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(fill)
    else:
        # Explicitly set "no fill" so the cell does NOT inherit banding
        # from the table style. Combined with _reset_table_style(), this
        # guarantees the slide background shows through cleanly.
        cell.fill.background()
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(40_000)
    tf.margin_right = Emu(40_000)
    tf.margin_top = Emu(20_000)
    tf.margin_bottom = Emu(20_000)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = sanitize_text(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


# ---------------------------------------------------------------------------
# Display helpers
# ---------------------------------------------------------------------------
def _display_name(entity: dict, snapshots: list | None = None) -> str:
    """Return the sanitized model/judge name.

    Previously this prefixed reasoning models with ``(R) `` because the
    exporter assumed PPTX had no inline SVG. The new pipeline renders a
    native brain-icon PNG via :func:`_add_brain_picture` at each call
    site, so the text is the plain name and the `snapshots` argument is
    kept only for signature compatibility with legacy callers.
    """
    return sanitize_text(entity.get("name", ""))


def _per_prompt_cost(model: dict, n_questions: int) -> str:
    cost = model.get("estimated_cost")
    if cost is None:
        return "—"
    return format_cost(round(cost / max(1, n_questions), 4))


def _format_score(score: float | None) -> str:
    if score is None:
        return "—"
    return f"{score:.2f}"


def _format_ci(ci: dict | None) -> str:
    if not ci:
        return ""
    return f"(95% CI {ci['lower']:.2f}–{ci['upper']:.2f})"


def _match_model_stats(name: str, statistics: dict | None) -> dict | None:
    if not statistics or not statistics.get("model_statistics"):
        return None
    for ms in statistics["model_statistics"]:
        if ms.get("model_name") == name:
            return ms
    return None


def _top_cluster_count(winner_ci: dict | None, statistics: dict | None) -> int | None:
    if not winner_ci or not statistics or not statistics.get("model_statistics"):
        return None
    lower = winner_ci.get("lower")
    if lower is None:
        return None
    count = 0
    for ms in statistics["model_statistics"]:
        ci = ms.get("weighted_score_ci")
        if ci and ci.get("upper") is not None and ci["upper"] >= lower:
            count += 1
    return count if count else None


def _export_date(run: dict) -> str:
    raw = run.get("completed_at") or run.get("created_at") or ""
    return sanitize_text(raw)[:10] if raw else ""


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def _blank_slide(prs: Presentation, tokens: BrandTokens):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, tokens)
    return slide


def _add_slide1_cover(
    prs: Presentation,
    data: dict,
    tokens: BrandTokens,
    theme: str,
) -> None:
    slide = _blank_slide(prs, tokens)
    run = data["run"]
    models = data["models"]
    questions = data["questions"]
    judges = data["judges"]
    criteria = run.get("criteria") or []
    snapshots = run.get("model_preset_snapshots")

    # Eyebrow + title
    _add_text(
        slide, "MODEL EVALUATION REPORT",
        x=80, y=70, w=1000, h=32,
        color=tokens["muted_foreground"], size=12, bold=True,
    )
    _add_text(
        slide, run.get("name", "Benchmark Run"),
        x=80, y=108, w=1760, h=86,
        color=tokens["foreground"], size=36, bold=True,
    )

    # KEY FINDING block (left column)
    winner = models[0]
    winner_name = _display_name(winner, snapshots)
    winner_score = _format_score(winner.get("weighted_score"))
    winner_stats = _match_model_stats(winner["name"], data.get("statistics"))
    winner_ci_text = _format_ci(
        winner_stats.get("weighted_score_ci") if winner_stats else None
    )

    _add_rect(
        slide, x=80, y=220, w=840, h=300,
        fill=tokens["card"], border=tokens["border"],
    )
    _add_text(
        slide, "KEY FINDING",
        x=104, y=244, w=800, h=28,
        color=tokens["muted_foreground"], size=12, bold=True,
    )
    # Winner name: amber-600 (light) / yellow-400 (dark) to match the
    # frontend winner card in OverviewSection.tsx. 2 lines at 22pt to
    # handle long labels like "Claude Opus 4.6 [Reasoning (high)]".
    _add_text(
        slide, winner_name,
        x=104, y=284, w=800, h=96,
        color=brand(theme, "winner_text"), size=22, bold=True,
    )
    _add_text(
        slide, f"Weighted mean score {winner_score} / 10",
        x=104, y=388, w=800, h=34,
        color=tokens["foreground"], size=16,
    )
    _add_text(
        slide,
        winner_ci_text or "(CI unavailable for this run)",
        x=104, y=424, w=800, h=28,
        color=tokens["muted_foreground"], size=13,
    )
    _add_text(
        slide,
        f"Top performer across {len(questions)} prompts judged by "
        f"{len(judges)} independent judges.",
        x=104, y=456, w=800, h=50,
        color=tokens["muted_foreground"], size=12,
    )

    # PODIUM block (right column)
    _add_rect(
        slide, x=960, y=220, w=880, h=300,
        fill=tokens["card"], border=tokens["border"],
    )
    _add_text(
        slide, "PODIUM",
        x=984, y=244, w=840, h=28,
        color=tokens["muted_foreground"], size=12, bold=True,
    )
    medals = ["1st", "2nd", "3rd"]
    for i, m in enumerate(models[:3]):
        row_y = 284 + i * 66
        _add_text(
            slide, medals[i],
            x=984, y=row_y, w=60, h=40,
            color=tokens["muted_foreground"], size=14, bold=True,
        )
        # Podium name color: rank 1 amber, 2 green, 3 orange (frontend parity)
        podium_key = "winner_text" if i == 0 else ("runner_text" if i == 1 else "bronze_text")
        _add_text(
            slide, _display_name(m, snapshots),
            x=1060, y=row_y, w=600, h=40,
            color=brand(theme, podium_key), size=16, bold=True,
        )
        _add_text(
            slide, _format_score(m.get("weighted_score")),
            x=1700, y=row_y, w=120, h=40,
            color=tokens["foreground"], size=16, bold=True,
            align=PP_ALIGN.RIGHT,
        )

    # Six scope tiles
    total_spend = format_cost(run.get("total_cost", 0.0))
    mode_str = (run.get("judge_mode") or "—").upper()
    scope_tiles: list[tuple[str, str]] = [
        ("CANDIDATES", str(len(models))),
        ("PROMPTS", str(len(questions))),
        ("JUDGES", str(len(judges))),
        ("RUBRIC CRITERIA", str(len(criteria))),
        ("TOTAL SPEND", total_spend),
        ("MODE", mode_str),
    ]
    tile_w = 280.0
    tile_h = 120.0
    tile_x0 = 80
    tile_y = 560
    gap = 12.0
    for i, (label, value) in enumerate(scope_tiles):
        col = i
        x = tile_x0 + col * (tile_w + gap)
        _add_stat_tile(
            slide, x=x, y=tile_y, w=tile_w, h=tile_h,
            label=label, value=value, tokens=tokens,
        )

    # Footer
    run_id = run.get("id") or 0
    _add_text(
        slide, f"BeLLMark Run #{run_id:04d}",
        x=80, y=1010, w=600, h=32,
        color=tokens["muted_foreground"], size=11,
    )
    _add_text(
        slide, "bellmark.ai",
        x=1240, y=1010, w=600, h=32,
        color=tokens["muted_foreground"], size=11,
        align=PP_ALIGN.RIGHT,
    )


def _add_slide2_executive(
    prs: Presentation,
    data: dict,
    tokens: BrandTokens,
    theme: str,
) -> None:
    slide = _blank_slide(prs, tokens)
    run = data["run"]
    models = data["models"]
    questions = data["questions"]
    snapshots = run.get("model_preset_snapshots")
    winner = models[0]
    statistics = data.get("statistics")
    winner_stats = _match_model_stats(winner["name"], statistics)
    winner_ci = winner_stats.get("weighted_score_ci") if winner_stats else None

    _add_text(
        slide, "EXECUTIVE SUMMARY",
        x=80, y=60, w=800, h=28,
        color=tokens["muted_foreground"], size=12, bold=True,
    )

    winner_name = _display_name(winner, snapshots)
    # Exec winner: amber/yellow text (frontend winner card parity).
    _add_text(
        slide, winner_name,
        x=80, y=96, w=1760, h=60,
        color=brand(theme, "winner_text"), size=30, bold=True,
    )
    score_text = f"{winner.get('weighted_score', 0):.2f} / 10"
    _add_text(
        slide, score_text,
        x=80, y=164, w=600, h=50,
        color=tokens["foreground"], size=24, bold=True,
    )
    ci_text = _format_ci(winner_ci) if winner_ci else "(CI unavailable for this run)"
    _add_text(
        slide, ci_text,
        x=440, y=172, w=800, h=40,
        color=tokens["muted_foreground"], size=14,
    )

    # Three stat tiles
    cluster_count = _top_cluster_count(winner_ci, statistics)
    top_cluster_val = f"{cluster_count} tied" if cluster_count else "—"
    tile_defs: list[tuple[str, str]] = [
        ("TOTAL SPEND", format_cost(run.get("total_cost", 0.0))),
        ("TOP CLUSTER", top_cluster_val),
        ("SAMPLE SIZE", f"{len(questions)} prompts"),
    ]
    tile_y = 250
    tile_w = 560.0
    tile_h = 140.0
    gap = 20.0
    for i, (label, value) in enumerate(tile_defs):
        x = 80 + i * (tile_w + gap)
        _add_stat_tile(
            slide, x=x, y=tile_y, w=tile_w, h=tile_h,
            label=label, value=value, tokens=tokens,
        )

    # Top 3 price/quality strip
    _add_text(
        slide, "TOP 3 · QUALITY VS COST PER PROMPT",
        x=80, y=440, w=1200, h=28,
        color=tokens["muted_foreground"], size=12, bold=True,
    )
    for i, m in enumerate(models[:3]):
        y = 488 + i * 124
        _add_rect(
            slide, x=80, y=y, w=1760, h=108,
            fill=tokens["card"], border=tokens["border"],
        )
        _add_text(
            slide, f"#{i + 1}",
            x=104, y=y + 18, w=80, h=40,
            color=tokens["muted_foreground"], size=18, bold=True,
        )
        # Top-3 band on executive slide — rank 1 amber, 2 green, 3 orange.
        band_key = "winner_text" if i == 0 else ("runner_text" if i == 1 else "bronze_text")
        _add_text(
            slide, _display_name(m, snapshots),
            x=200, y=y + 18, w=900, h=40,
            color=brand(theme, band_key), size=18, bold=True,
        )
        _add_text(
            slide, f"{m.get('weighted_score', 0):.2f} / 10",
            x=200, y=y + 60, w=400, h=30,
            color=tokens["muted_foreground"], size=13,
        )
        per_prompt = _per_prompt_cost(m, len(questions))
        marker = per_prompt if per_prompt == "—" else f"{per_prompt}/prompt"
        _add_text(
            slide, marker,
            x=1520, y=y + 30, w=220, h=48,
            color=tokens["foreground"], size=20, bold=True,
            align=PP_ALIGN.RIGHT,
        )


def _add_slide3_leaderboard(
    prs: Presentation,
    data: dict,
    tokens: BrandTokens,
    theme: str,
) -> None:
    slide = _blank_slide(prs, tokens)
    run = data["run"]
    models = data["models"]
    questions = data["questions"]
    statistics = data.get("statistics")
    snapshots = run.get("model_preset_snapshots")
    n_q = len(questions)

    _add_text(
        slide, "LEADERBOARD",
        x=80, y=60, w=800, h=28,
        color=tokens["muted_foreground"], size=12, bold=True,
    )
    _add_text(
        slide, "Ranking with uncertainty, speed, win rates, and per-prompt cost",
        x=80, y=96, w=1760, h=48,
        color=tokens["foreground"], size=24, bold=True,
    )

    headers = [
        "#", "MODEL", "MEAN", "95% CI", "WIN RATE",
        "TOK/S", "AVG LATENCY", "LC WIN", "$/PROMPT",
    ]
    rows = 1 + len(models)
    cols = len(headers)
    table_x = 80
    table_y = 180
    table_w = 1728  # normative per spec §2.2 — matches cross-format width
    table_h = min(820, 60 + 64 * len(models))
    table_shape = slide.shapes.add_table(
        rows, cols, lx(table_x), ly(table_y), lx(table_w), ly(table_h)
    )
    table = table_shape.table
    _reset_table_style(table)

    # Percent widths per spec §2.2 (sum = 100). Logical-px widths are
    # computed from 1728px table width; use `lx(...)` for EMU.
    pcts = [4, 28, 7, 13, 9, 8, 11, 9, 11]
    col_widths_px = [round(table_w * p / 100) for p in pcts]
    for i, w in enumerate(col_widths_px):
        table.columns[i].width = lx(w)

    # Header row
    for i, h in enumerate(headers):
        _set_cell(
            table.cell(0, i), h,
            color=tokens["foreground"], size=12, bold=True,
            fill=tokens["muted"],
            align=PP_ALIGN.LEFT if i != 0 else PP_ALIGN.CENTER,
        )

    # Row height (logical-px) — used to anchor brain icons inside cells.
    row_h_px = (table_h - 60) / max(1, len(models))

    for row_idx, m in enumerate(models, start=1):
        ms = _match_model_stats(m["name"], statistics)

        # # column
        _set_cell(
            table.cell(row_idx, 0), f"#{m.get('rank') or row_idx}",
            color=tokens["muted_foreground"], size=12, bold=True,
            align=PP_ALIGN.CENTER,
        )

        # MODEL — winner gets amber/yellow brand text, others green (mirrors
        # frontend `text-amber-600 dark:text-yellow-400` / `text-green-600
        # dark:text-green-400`). Reasoning marker is the "[Reasoning (high)]"
        # suffix already inside the name; no icon needed.
        model_color = brand(theme, "winner_text" if row_idx == 1 else "runner_text")
        _set_cell(
            table.cell(row_idx, 1), _display_name(m, snapshots),
            color=model_color, size=12, bold=(row_idx == 1),
        )

        # MEAN
        mean = m.get("weighted_score")
        _set_cell(
            table.cell(row_idx, 2),
            f"{float(mean):.2f}" if mean is not None else "—",
            color=tokens["foreground"], size=12,
        )

        # 95% CI
        ci = ms.get("weighted_score_ci") if ms else None
        ci_cell = (
            f"[{float(ci['lower']):.2f} – {float(ci['upper']):.2f}]"
            if ci and ci.get("lower") is not None and ci.get("upper") is not None
            else "—"
        )
        _set_cell(
            table.cell(row_idx, 3), ci_cell,
            color=tokens["muted_foreground"], size=11,
        )

        # WIN RATE (point estimate only)
        wr = m.get("win_rate")
        wr_cell = f"{float(wr) * 100:.0f}%" if wr is not None else "—"
        _set_cell(
            table.cell(row_idx, 4), wr_cell,
            color=tokens["foreground"], size=11,
        )

        # TOK/S
        tps = m.get("tokens_per_second")
        tps_cell = f"{float(tps):.1f}" if tps is not None else "—"
        _set_cell(
            table.cell(row_idx, 5), tps_cell,
            color=tokens["foreground"], size=11,
        )

        # AVG LATENCY
        _set_cell(
            table.cell(row_idx, 6), _fmt_latency_ms(m.get("avg_latency_ms")),
            color=tokens["foreground"], size=11,
        )

        # LC WIN (with length-bias warning)
        lc_blk = (ms or {}).get("lc_win_rate") or {}
        lc_val = lc_blk.get("lc_win_rate")
        if lc_val is not None:
            warn = "⚠ " if lc_blk.get("length_bias_detected") else ""
            lc_cell = f"{warn}{float(lc_val) * 100:.0f}%"
        else:
            lc_cell = "—"
        _set_cell(
            table.cell(row_idx, 7), lc_cell,
            color=tokens["foreground"], size=11,
        )

        # $/PROMPT
        cost = m.get("estimated_cost")
        if cost is not None and n_q > 0:
            cost_cell = f"{format_cost(cost / n_q)}/prompt"
        else:
            cost_cell = "—"
        _set_cell(
            table.cell(row_idx, 8), cost_cell,
            color=tokens["foreground"], size=11,
        )


def _score_cell_fill(score: float, tokens: BrandTokens) -> tuple[int, int, int]:
    """Blend muted→chart_2 (or chart_1 in dark) based on 0..10 score."""
    base = tokens["muted"]
    accent = tokens["chart_2"] if max(tokens["chart_2"]) > 80 else tokens["chart_1"]
    # Clamp
    s = max(0.0, min(10.0, float(score))) / 10.0
    # Light blend: keep cells legible (cap accent influence at 45%).
    t = 0.15 + 0.3 * s
    return (
        round(base[0] * (1 - t) + accent[0] * t),
        round(base[1] * (1 - t) + accent[1] * t),
        round(base[2] * (1 - t) + accent[2] * t),
    )


def _add_slide4_criteria(
    prs: Presentation,
    data: dict,
    tokens: BrandTokens,
    theme: str,
) -> None:
    slide = _blank_slide(prs, tokens)
    run = data["run"]
    models = data["models"]
    criteria = run.get("criteria") or []
    snapshots = run.get("model_preset_snapshots")
    crit_names = [c["name"] for c in criteria]

    _add_text(
        slide, "PER-CRITERION SCORES",
        x=80, y=60, w=800, h=28,
        color=tokens["muted_foreground"], size=12, bold=True,
    )
    _add_text(
        slide, "How each model performed against every rubric criterion",
        x=80, y=96, w=1760, h=48,
        color=tokens["foreground"], size=24, bold=True,
    )

    if not crit_names:
        _add_text(
            slide, "No criteria recorded for this run.",
            x=80, y=240, w=1760, h=80,
            color=tokens["muted_foreground"], size=16,
        )
        return

    headers = ["Model"] + crit_names
    rows = 1 + len(models)
    cols = len(headers)
    table_x = 80
    table_y = 180
    table_w = 1760
    table_h = min(820, 70 + 64 * len(models))
    table_shape = slide.shapes.add_table(
        rows, cols, lx(table_x), ly(table_y), lx(table_w), ly(table_h)
    )
    table = table_shape.table
    _reset_table_style(table)

    # First column wider for model names.
    name_col_w = 560
    per_crit_w = (1760 - name_col_w) // max(1, len(crit_names))
    table.columns[0].width = lx(name_col_w)
    for i in range(1, cols):
        table.columns[i].width = lx(per_crit_w)

    # Header row
    _set_cell(
        table.cell(0, 0), "Model",
        color=tokens["foreground"], size=12, bold=True,
        fill=tokens["muted"],
    )
    for i, name in enumerate(crit_names, start=1):
        # Criteria header: amber text + amber fill (mirrors criteria badge).
        _set_cell(
            table.cell(0, i), name,
            color=brand(theme, "criteria_text"), size=10, bold=True,
            fill=brand(theme, "criteria_bg"), align=PP_ALIGN.CENTER,
        )

    # Best-in-class tracking
    best_by_crit: dict[str, tuple[str, float]] = {}
    for crit in crit_names:
        best_model = None
        best_score = -1.0
        for m in models:
            score = (m.get("per_criterion_scores") or {}).get(crit)
            if score is not None and score > best_score:
                best_score = score
                best_model = m["name"]
        if best_model:
            best_by_crit[crit] = (best_model, best_score)

    # Row height (logical-px) for brain-icon anchor math on the criteria table.
    for row_idx, m in enumerate(models, start=1):
        name_color = brand(theme, "winner_text" if row_idx == 1 else "runner_text")
        _set_cell(
            table.cell(row_idx, 0), _display_name(m, snapshots),
            color=name_color, size=11, bold=(row_idx == 1),
        )
        scores = m.get("per_criterion_scores") or {}
        for col_idx, crit in enumerate(crit_names, start=1):
            s = scores.get(crit)
            text = f"{s:.1f}" if s is not None else "—"
            fill = _score_cell_fill(s if s is not None else 0.0, tokens)
            _set_cell(
                table.cell(row_idx, col_idx), text,
                color=tokens["foreground"], size=11,
                fill=fill, align=PP_ALIGN.CENTER,
            )

    # Best-in-class strip below
    strip_y = table_y + table_h + 30
    _add_text(
        slide, "BEST-IN-CLASS PER CRITERION",
        x=80, y=strip_y, w=1760, h=24,
        color=tokens["muted_foreground"], size=11, bold=True,
    )
    line = " · ".join(
        f"{crit}: {name} ({score:.1f})"
        for crit, (name, score) in best_by_crit.items()
    )
    _add_text(
        slide, line or "—",
        x=80, y=strip_y + 30, w=1760, h=60,
        color=tokens["foreground"], size=12,
    )


def _add_slide5_stats_rigor(
    prs: Presentation,
    data: dict,
    tokens: BrandTokens,
    theme: str,
) -> None:
    slide = _blank_slide(prs, tokens)
    run = data["run"]
    models = data["models"]
    statistics = data.get("statistics")
    snapshots = run.get("model_preset_snapshots")

    _add_text(
        slide, "STATISTICAL RIGOR",
        x=80, y=60, w=800, h=28,
        color=tokens["muted_foreground"], size=12, bold=True,
    )
    _add_text(
        slide, "Methodology, effect sizes, and pairwise comparisons",
        x=80, y=96, w=1760, h=48,
        color=tokens["foreground"], size=24, bold=True,
    )

    # Static methodology paragraph — ALWAYS present, contains all 9 terms.
    _add_rect(
        slide, x=80, y=180, w=1760, h=200,
        fill=tokens["card"], border=tokens["border"],
    )
    _add_text(
        slide, "METHODOLOGY",
        x=104, y=200, w=1720, h=24,
        color=tokens["muted_foreground"], size=11, bold=True,
    )
    _add_text(
        slide, _METHODOLOGY_PARAGRAPH,
        x=104, y=226, w=1720, h=150,
        color=tokens["foreground"], size=12,
    )

    # Pairwise matrix area
    n = min(8, len(models))
    m_codes = [f"M{i + 1}" for i in range(n)]
    model_names_by_code = {m_codes[i]: models[i]["name"] for i in range(n)}

    matrix_y = 400
    matrix_h = 420

    pairwise = None
    if statistics:
        pairwise = statistics.get("pairwise_comparisons")

    if not pairwise:
        _add_rect(
            slide, x=80, y=matrix_y, w=1160, h=matrix_h,
            fill=tokens["card"], border=tokens["border"],
        )
        _add_text(
            slide, "Statistical analysis unavailable for this run.",
            x=104, y=matrix_y + 40, w=1120, h=80,
            color=tokens["muted_foreground"], size=16,
        )
    else:
        # Native table: header row is M-codes, first col is M-codes.
        rows = n + 1
        cols = n + 1
        table_shape = slide.shapes.add_table(
            rows, cols, lx(80), ly(matrix_y), lx(1160), ly(matrix_h)
        )
        table = table_shape.table
        _reset_table_style(table)
        # Header corner + column headers
        _set_cell(
            table.cell(0, 0), "",
            color=tokens["foreground"], fill=tokens["muted"], size=11,
        )
        for j in range(n):
            _set_cell(
                table.cell(0, j + 1), m_codes[j],
                color=tokens["foreground"], size=11, bold=True,
                fill=tokens["muted"], align=PP_ALIGN.CENTER,
            )
        # Body
        pair_index: dict[frozenset, dict] = {}
        for pc in pairwise:
            key = frozenset({pc.get("model_a"), pc.get("model_b")})
            pair_index[key] = pc

        for i in range(n):
            _set_cell(
                table.cell(i + 1, 0), m_codes[i],
                color=tokens["foreground"], size=11, bold=True,
                fill=tokens["muted"], align=PP_ALIGN.CENTER,
            )
            for j in range(n):
                if i == j:
                    _set_cell(
                        table.cell(i + 1, j + 1), "—",
                        color=tokens["muted_foreground"], size=10,
                        align=PP_ALIGN.CENTER,
                    )
                    continue
                key = frozenset({models[i].get("name", ""), models[j].get("name", "")})
                pc = pair_index.get(key)
                if not pc:
                    text = ""
                else:
                    d = pc.get("cohens_d")
                    ap = pc.get("adjusted_p")
                    d_text = f"d={d:.2f}" if d is not None else "d=—"
                    ap_text = f"p={ap:.3f}" if ap is not None else "p=—"
                    text = f"{d_text}\n{ap_text}"
                _set_cell(
                    table.cell(i + 1, j + 1), text,
                    color=tokens["foreground"], size=9,
                    align=PP_ALIGN.CENTER,
                )

    # M-code legend (right column)
    legend_x = 1280
    legend_y = 400
    _add_text(
        slide, "M-CODE LEGEND",
        x=legend_x, y=legend_y, w=560, h=24,
        color=tokens["muted_foreground"], size=11, bold=True,
    )
    for i, code in enumerate(m_codes):
        name = _display_name(models[i], snapshots)
        y_pos = legend_y + 28 + i * 22
        # M-code legend colors: rank 1 amber, others green.
        legend_color = brand(theme, "winner_text" if i == 0 else "runner_text")
        _add_text(
            slide, f"{code} = {name}",
            x=legend_x, y=y_pos, w=560, h=22,
            color=legend_color, size=11,
        )

    # HOW TO READ + EFFECT SIZE legends
    legend2_y = legend_y + 28 + n * 22 + 20
    _add_text(
        slide, "HOW TO READ",
        x=legend_x, y=legend2_y, w=560, h=20,
        color=tokens["muted_foreground"], size=10, bold=True,
    )
    _add_text(
        slide, "Row ahead of column · Row behind column · No significant difference",
        x=legend_x, y=legend2_y + 22, w=560, h=60,
        color=tokens["foreground"], size=10,
    )
    _add_text(
        slide, "EFFECT SIZE (COHEN'S D)",
        x=80, y=matrix_y + matrix_h + 10, w=560, h=20,
        color=tokens["muted_foreground"], size=10, bold=True,
    )
    _add_text(
        slide,
        "Negligible (|d| < 0.2) · Small (|d| < 0.5) · Medium (|d| < 0.8) · Large (|d| >= 0.8)",
        x=80, y=matrix_y + matrix_h + 32, w=1760, h=40,
        color=tokens["foreground"], size=10,
    )


def _bias_panel(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    panel: dict | None,
    tokens: BrandTokens,
) -> None:
    _add_rect(
        slide, x=x, y=y, w=w, h=h,
        fill=tokens["card"], border=tokens["border"],
    )
    _add_text(
        slide, title,
        x=x + 20, y=y + 18, w=w - 40, h=30,
        color=tokens["foreground"], size=16, bold=True,
    )
    if not panel:
        _add_text(
            slide, "Not available for this run.",
            x=x + 20, y=y + 56, w=w - 40, h=40,
            color=tokens["muted_foreground"], size=12,
        )
        return

    severity = (panel.get("severity") or "none").lower()
    detected = severity not in ("none", "low")
    status = "Detected" if detected else "Not detected"
    status_color = tokens["destructive"] if detected else tokens["muted_foreground"]
    _add_text(
        slide, status,
        x=x + 20, y=y + 52, w=w - 40, h=26,
        color=status_color, size=12, bold=True,
    )
    corr = panel.get("correlation")
    p = panel.get("p_value")
    stat_text_parts = []
    if corr is not None:
        stat_text_parts.append(f"r = {corr:.3f}")
    if p is not None:
        stat_text_parts.append(f"p = {p:.3f}")
    if stat_text_parts:
        _add_text(
            slide, " · ".join(stat_text_parts),
            x=x + 20, y=y + 82, w=w - 40, h=24,
            color=tokens["foreground"], size=11,
        )
    desc = panel.get("description") or ""
    _add_text(
        slide, desc,
        x=x + 20, y=y + 110, w=w - 40, h=h - 120,
        color=tokens["muted_foreground"], size=10,
    )


def _add_slide6_bias(
    prs: Presentation,
    data: dict,
    tokens: BrandTokens,
    theme: str,
) -> None:
    slide = _blank_slide(prs, tokens)
    bias = data.get("bias_report")
    judge_summary = data.get("judge_summary") or {}
    kappa_value = data.get("kappa_value")
    kappa_type = data.get("kappa_type")

    _add_text(
        slide, "BIAS & CALIBRATION",
        x=80, y=60, w=800, h=28,
        color=tokens["muted_foreground"], size=12, bold=True,
    )
    _add_text(
        slide, "Bias screens and judge calibration",
        x=80, y=96, w=1760, h=48,
        color=tokens["foreground"], size=24, bold=True,
    )

    # 2x2 bias grid
    grid_x = 80
    grid_y = 180
    panel_w = 870.0
    panel_h = 240.0
    gap = 20.0
    panels_spec: list[tuple[str, str]] = [
        ("Position bias", "position_bias"),
        ("Length bias", "length_bias"),
        ("Self-preference", "self_preference"),
        ("Verbosity bias", "verbosity_bias"),
    ]
    if not bias:
        _add_rect(
            slide, x=grid_x, y=grid_y, w=1760, h=500,
            fill=tokens["card"], border=tokens["border"],
        )
        _add_text(
            slide, "Bias analysis unavailable for this run.",
            x=grid_x + 30, y=grid_y + 40, w=1700, h=60,
            color=tokens["muted_foreground"], size=18,
        )
        # Still mention all four labels so Scene 4b term search passes.
        label_str = " · ".join(label for label, _ in panels_spec)
        _add_text(
            slide, label_str,
            x=grid_x + 30, y=grid_y + 120, w=1700, h=40,
            color=tokens["muted_foreground"], size=14,
        )
    else:
        for i, (label, key) in enumerate(panels_spec):
            col = i % 2
            row = i // 2
            x = grid_x + col * (panel_w + gap)
            y = grid_y + row * (panel_h + gap)
            _bias_panel(
                slide, x=x, y=y, w=panel_w, h=panel_h,
                title=label, panel=bias.get(key), tokens=tokens,
            )

    # Calibration strip (always rendered from judge_summary + kappa_value)
    strip_y = 700
    tile_w = 560.0
    tile_h = 160.0
    tile_gap = 20.0

    agreement = judge_summary.get("agreement_rate")
    agreement_val = f"{agreement * 100:.0f}%" if agreement is not None else "—"

    kappa_label = "FLEISS' K" if kappa_type == "fleiss" else "COHEN'S K"
    if kappa_type is None:
        kappa_label = "COHEN'S K"
    kappa_val = f"{kappa_value:.3f}" if isinstance(kappa_value, (int, float)) else "—"

    disagree = judge_summary.get("disagreement_count")
    disagree_val = str(disagree) if disagree is not None else "—"

    cal_tiles = [
        ("INTER-JUDGE AGREEMENT", agreement_val),
        (kappa_label, kappa_val),
        ("DISAGREEMENT QUESTIONS", disagree_val),
    ]
    for i, (label, value) in enumerate(cal_tiles):
        x = 80 + i * (tile_w + tile_gap)
        _add_stat_tile(
            slide, x=x, y=strip_y, w=tile_w, h=tile_h,
            label=label, value=value, tokens=tokens,
        )


def _add_slide7_methodology(
    prs: Presentation,
    data: dict,
    tokens: BrandTokens,
    theme: str,
) -> None:
    slide = _blank_slide(prs, tokens)
    run = data["run"]
    models = data["models"]
    judges = data["judges"]
    questions = data["questions"]
    criteria = run.get("criteria") or []
    snapshots = run.get("judge_preset_snapshots")

    _add_text(
        slide, "METHODOLOGY & SIGN-OFF",
        x=80, y=60, w=800, h=28,
        color=tokens["muted_foreground"], size=12, bold=True,
    )
    _add_text(
        slide, "How this report was produced",
        x=80, y=96, w=1760, h=48,
        color=tokens["foreground"], size=24, bold=True,
    )

    # SCOPE section
    _add_text(
        slide, "SCOPE",
        x=80, y=180, w=860, h=24,
        color=tokens["muted_foreground"], size=11, bold=True,
    )
    judgments_total = len(questions) * max(1, len(judges))
    scope_items = [
        ("Prompts", str(len(questions))),
        ("Candidate models", str(len(models))),
        ("Judges", str(len(judges))),
        ("Rubric criteria", str(len(criteria))),
        ("Judging mode", run.get("judge_mode") or "—"),
        ("Judgments total", str(judgments_total)),
        ("Total spend", format_cost(run.get("total_cost", 0.0))),
    ]
    for i, (label, value) in enumerate(scope_items):
        _add_text(
            slide, label,
            x=80, y=210 + i * 30, w=400, h=26,
            color=tokens["muted_foreground"], size=12,
        )
        _add_text(
            slide, value,
            x=500, y=210 + i * 30, w=440, h=26,
            color=tokens["foreground"], size=12, bold=True,
        )

    # RUBRIC CRITERIA section
    _add_text(
        slide, "RUBRIC CRITERIA",
        x=980, y=180, w=860, h=24,
        color=tokens["muted_foreground"], size=11, bold=True,
    )
    if criteria:
        for i, c in enumerate(criteria):
            line = f"• {c['name']} (weight {c.get('weight', 1.0)})"
            _add_text(
                slide, line,
                x=980, y=210 + i * 30, w=860, h=28,
                color=tokens["foreground"], size=12,
            )
    else:
        _add_text(
            slide, "No rubric criteria recorded.",
            x=980, y=210, w=860, h=28,
            color=tokens["muted_foreground"], size=12,
        )

    # JUDGE PANEL section
    judge_y = 460
    _add_text(
        slide, "JUDGE PANEL",
        x=80, y=judge_y, w=860, h=24,
        color=tokens["muted_foreground"], size=11, bold=True,
    )
    if judges:
        for i, j in enumerate(judges):
            line_y = judge_y + 30 + i * 28
            # Judges are purple in the frontend (judge badge class). We render
            # the bullet in muted-foreground and the name in purple.
            _add_text(
                slide, "•",
                x=80, y=line_y, w=20, h=28,
                color=tokens["muted_foreground"], size=12,
            )
            _add_text(
                slide, _display_name(j, snapshots),
                x=104, y=line_y, w=840, h=28,
                color=brand(theme, "judge_text"), size=12,
            )
    else:
        _add_text(
            slide, "No judges recorded.",
            x=80, y=judge_y + 30, w=860, h=28,
            color=tokens["muted_foreground"], size=12,
        )

    # SIGN-OFF section
    _add_text(
        slide, "SIGN-OFF",
        x=980, y=judge_y, w=860, h=24,
        color=tokens["muted_foreground"], size=11, bold=True,
    )
    export_date = _export_date(run)
    signoff_items = [
        ("Prepared by", "BeLLMark Studio"),
        ("Reviewed by", "—"),
        ("Export date", export_date or "—"),
        ("Source", f"BeLLMark Run #{run.get('id', 0):04d}"),
        ("Methodology", "See slide 5 — Wilson / Bootstrap / Wilcoxon / Holm-Bonferroni"),
    ]
    for i, (label, value) in enumerate(signoff_items):
        _add_text(
            slide, label,
            x=980, y=judge_y + 30 + i * 32, w=300, h=26,
            color=tokens["muted_foreground"], size=12,
        )
        _add_text(
            slide, value,
            x=1300, y=judge_y + 30 + i * 32, w=540, h=26,
            color=tokens["foreground"], size=12, bold=True,
        )

    # Footer
    _add_text(
        slide, "BeLLMark · bellmark.ai",
        x=80, y=1010, w=600, h=28,
        color=tokens["muted_foreground"], size=11,
    )


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------
def generate_pptx(data: dict, theme: str = "light") -> bytes:
    """Render the 7-slide McKinsey-grade PPTX.

    Args:
        data: Output of :func:`app.core.exports.common.prepare_export_data`.
        theme: ``"light"`` or ``"dark"``. Controls the brand token palette.

    Returns:
        Bytes of a valid .pptx file. Always 7 slides, regardless of whether
        optional analytics blocks (``statistics``, ``bias_report``,
        ``calibration_report``) are populated.
    """
    theme_name = theme if theme in ("light", "dark") else "light"
    tokens = get_tokens(theme_name)

    prs = Presentation()
    from pptx.util import Inches

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_slide1_cover(prs, data, tokens, theme_name)
    _add_slide2_executive(prs, data, tokens, theme_name)
    _add_slide3_leaderboard(prs, data, tokens, theme_name)
    _add_slide4_criteria(prs, data, tokens, theme_name)
    _add_slide5_stats_rigor(prs, data, tokens, theme_name)
    _add_slide6_bias(prs, data, tokens, theme_name)
    _add_slide7_methodology(prs, data, tokens, theme_name)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


__all__ = ["generate_pptx"]
